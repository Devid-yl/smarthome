"""
Générateur de présentation PowerPoint pour le projet SmartHome
Durée: 5-10 minutes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ==================== SLIDE 1: TITRE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Titre principal
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "SmartHome"
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(37, 99, 235)  # Bleu
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Application de Domotique Intelligente"
    
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(100, 116, 139)
    p.alignment = PP_ALIGN.CENTER
    
    # Infos
    info_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1))
    info_frame = info_box.text_frame
    info_frame.text = "Python Tornado • PostgreSQL • WebSocket\nArchitecture REST API complète"
    
    for p in info_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(148, 163, 184)
        p.alignment = PP_ALIGN.CENTER
    
    # ==================== SLIDE 2: ARCHITECTURE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    title = slide.shapes.title
    title.text = "Architecture Technique"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    # Contenu
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """Stack Technologique:
    
Backend:
• Python Tornado 6.4+ (async/await natif)
• SQLAlchemy 2.0 (ORM asynchrone)
• PostgreSQL 15+ avec asyncpg
• bcrypt pour hachage sécurisé

Frontend:
• HTML5 / CSS3 (Grid & Flexbox)
• JavaScript ES6+ Vanilla
• WebSocket temps réel
• Architecture SPA (Single Page Application)

Infrastructure:
• 9 tables PostgreSQL normalisées
• 50+ endpoints REST API
• WebSocket pour synchronisation temps réel"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.space_after = Pt(8)
        if paragraph.text.endswith(':'):
            paragraph.font.bold = True
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(37, 99, 235)
        elif paragraph.text.startswith('•'):
            paragraph.level = 1
    
    # ==================== SLIDE 3: STRUCTURE CODE ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Structure du Code"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """smarthome/
├── smarthome/tornado_app/
│   ├── models.py              # 9 modèles SQLAlchemy (User, House, Sensor...)
│   ├── database.py            # Configuration PostgreSQL async
│   ├── auth.py                # Authentification bcrypt + JWT
│   ├── app.py                 # Serveur Tornado + Routes (50+ endpoints)
│   │
│   ├── handlers/              # 15 handlers REST API
│   │   ├── users_api.py       # Auth, profil, upload photo
│   │   ├── houses_api.py      # CRUD maisons/pièces
│   │   ├── sensors.py         # 4 types capteurs IoT
│   │   ├── equipments.py      # 4 types équipements
│   │   ├── automation_rules.py # Règles conditionnelles
│   │   ├── house_members.py   # Invitations, rôles
│   │   ├── event_history.py   # Journalisation complète
│   │   └── websocket.py       # Temps réel
│   │
│   ├── services/
│   │   └── weather_service.py # API Open-Meteo
│   │
│   └── utils/
│       ├── grid_layers.py     # Système grille maison
│       └── permissions.py     # Gestion permissions
│
├── static/app/                # Frontend SPA
│   ├── dashboard.html         # Liste maisons
│   ├── house.html             # Détails maison + contrôles
│   ├── house.js               # Logique (2200+ lignes)
│   └── profile.html           # Profil utilisateur"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.name = 'Consolas'
        paragraph.space_after = Pt(4)
    
    # ==================== SLIDE 4: BACKEND - MODÈLES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Backend: Modèles de Données"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """# models.py - SQLAlchemy 2.0 (Async)

class User(Base):
    __tablename__ = 'users'
    id = Column(Integer, primary_key=True)
    username = Column(String(80), unique=True, nullable=False)
    email = Column(String(120), unique=True, nullable=False)
    password_hash = Column(String(255))  # bcrypt
    phone_number = Column(String(20))
    profile_image = Column(String(255))
    created_at = Column(DateTime, default=datetime.utcnow)

class Sensor(Base):
    __tablename__ = 'sensors'
    id = Column(Integer, primary_key=True)
    house_id = Column(Integer, ForeignKey('houses.id', ondelete='CASCADE'))
    name = Column(String(100), nullable=False)
    type = Column(String(50))  # temperature, luminosity, rain, presence
    value = Column(Float, default=0.0)
    unit = Column(String(20))  # °C, lux, %, boolean
    is_active = Column(Boolean, default=True)

class AutomationRule(Base):
    __tablename__ = 'automation_rules'
    id = Column(Integer, primary_key=True)
    sensor_id = Column(Integer, ForeignKey('sensors.id'))
    condition_operator = Column(String(10))  # >, <, >=, <=, ==, !=
    condition_value = Column(Float)
    equipment_id = Column(Integer, ForeignKey('equipments.id'))
    action_state = Column(String(20))  # on, off, open, closed
    is_active = Column(Boolean, default=True)"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.font.name = 'Consolas'
        paragraph.space_after = Pt(2)
        if paragraph.text.startswith('#'):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(34, 197, 94)
    
    # ==================== SLIDE 5: BACKEND - API REST ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Backend: API REST (50+ Endpoints)"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """# handlers/users_api.py - Authentification

class LoginAPIHandler(BaseAPIHandler):
    async def post(self):
        data = json.loads(self.request.body)
        username = data.get('username')
        password = data.get('password')
        
        async with async_session_maker() as session:
            result = await session.execute(
                select(User).where(User.username == username)
            )
            user = result.scalar_one_or_none()
            
            if user and verify_password(password, user.password_hash):
                self.set_secure_cookie("uid", str(user.id))
                self.write_json({"user": {...}})

# handlers/automation.py - Règles B2B

class AutomationTriggerHandler(BaseAPIHandler):
    async def post(self):
        async with async_session_maker() as session:
            # Récupérer règles actives
            rules = await session.execute(
                select(AutomationRule).where(
                    AutomationRule.is_active == True
                )
            )
            
            for rule in rules.scalars():
                # Évaluer condition
                if self._evaluate_condition(sensor.value, 
                    rule.condition_operator, rule.condition_value):
                    # Exécuter action
                    equipment.state = rule.action_state
                    await session.commit()"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.font.name = 'Consolas'
        paragraph.space_after = Pt(2)
        if paragraph.text.startswith('#'):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(34, 197, 94)
    
    # ==================== SLIDE 6: WEBSOCKET ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "WebSocket Temps Réel"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """# handlers/websocket.py - Synchronisation temps réel

class RealtimeHandler(tornado.websocket.WebSocketHandler):
    clients: Set["RealtimeHandler"] = set()
    
    def open(self):
        user_id = self.get_current_user()
        RealtimeHandler.clients.add(self)
        print(f"Client connecté: {user_id}, Total: {len(self.clients)}")
    
    @classmethod
    def broadcast_sensor_update(cls, sensor_id, value, house_id):
        message = json.dumps({
            "type": "sensor_update",
            "house_id": house_id,
            "data": {"id": sensor_id, "value": value}
        })
        
        for client in cls.clients:
            try:
                client.write_message(message)
            except Exception as e:
                print(f"Erreur broadcast: {e}")

# Frontend: realtime.js - Client WebSocket

function connectWebSocket() {
    const ws = new WebSocket(`ws://localhost:8001/ws/${houseId}`);
    
    ws.onmessage = (event) => {
        const data = JSON.parse(event.data);
        
        if (data.type === 'sensor_update') {
            updateSensorDisplay(data.data);
        } else if (data.type === 'equipment_update') {
            updateEquipmentDisplay(data.data);
        }
    };
}"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.font.name = 'Consolas'
        paragraph.space_after = Pt(2)
        if paragraph.text.startswith('#'):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(34, 197, 94)
    
    # ==================== SLIDE 7: FRONTEND ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Frontend: JavaScript (house.js)"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """// house.js - Logique page maison (2200+ lignes)

// Chargement capteurs avec tri stable
async function loadSensors() {
    const response = await fetch('/api/sensors');
    const data = await response.json();
    
    sensors = data.sensors
        .filter(s => s.house_id === parseInt(houseId))
        .sort((a, b) => a.id - b.id);  // Ordre stable
    
    displaySensors();
}

// Affichage grille interactive avec icônes
function displayHouseGrid() {
    const grid = house.grid;
    let html = '<table>';
    
    for (let i = 0; i < grid.length; i++) {
        html += '<tr>';
        for (let j = 0; j < grid[i].length; j++) {
            const cell = grid[i][j];
            const baseValue = getCellBase(cell);
            const sensors = getCellSensors(cell);
            const equipments = getCellEquipments(cell);
            
            // Afficher pièce + capteurs + équipements + utilisateurs
            html += `<td onclick="handleCellClick(${j}, ${i})">
                ${renderRoomName(baseValue)}
                ${renderSensorIcons(sensors)}
                ${renderEquipmentIcons(equipments)}
                ${renderUserAvatars(i, j)}
            </td>`;
        }
    }
    gridContainer.innerHTML = html;
}"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(11)
        paragraph.font.name = 'Consolas'
        paragraph.space_after = Pt(2)
        if paragraph.text.startswith('//'):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(34, 197, 94)
    
    # ==================== SLIDE 8: BASE DE DONNÉES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Base de Données PostgreSQL"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """9 Tables Normalisées:

users (1) ─── (N) houses (1) ─┬─ (N) rooms
                               ├─ (N) sensors
                               ├─ (N) equipments
                               ├─ (N) automation_rules
                               ├─ (N) house_members
                               ├─ (N) event_history
                               └─ (N) user_positions

Requête Complexe (Jointures multiples):

SELECT 
    h.name AS house,
    COUNT(DISTINCT s.id) AS sensors,
    COUNT(DISTINCT e.id) AS equipments,
    COUNT(DISTINCT ar.id) AS rules,
    COUNT(DISTINCT eh.id) AS events
FROM houses h
LEFT JOIN sensors s ON s.house_id = h.id
LEFT JOIN equipments e ON e.house_id = h.id
LEFT JOIN automation_rules ar ON ar.house_id = h.id
LEFT JOIN event_history eh ON eh.house_id = h.id
WHERE h.id = 1
GROUP BY h.id, h.name;

Contraintes d'Intégrité:
• Foreign Keys avec CASCADE DELETE
• Indexes sur colonnes fréquentes
• JSONB pour grilles et metadata
• Timestamps automatiques"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(12)
        if paragraph.text.endswith(':'):
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(37, 99, 235)
        elif 'SELECT' in paragraph.text or 'FROM' in paragraph.text:
            paragraph.font.name = 'Consolas'
            paragraph.font.size = Pt(10)
    
    # ==================== SLIDE 9: FONCTIONNALITÉS CLÉS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Fonctionnalités Principales"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """Authentification & Sécurité:
• Hachage bcrypt avec salt automatique
• Cookies HTTPOnly sécurisés
• JWT optionnel pour API REST
• Upload images avec validation (5MB max)

IoT & Automatisation:
• 4 types capteurs: Température, Luminosité, Pluie, Présence
• 4 types équipements: Volets, Portes, Lumières, Système sonore
• Règles conditionnelles (6 opérateurs: >, <, >=, <=, ==, !=)
• Déclenchement automatique basé capteurs

Multi-utilisateurs:
• Système d'invitations avec statut (pending/accepted/rejected)
• 2 rôles: Administrateur, Occupant
• Permissions granulaires par équipement
• Historique complet des actions

Temps Réel:
• WebSocket pour synchronisation instantanée
• 6 types de messages (sensor_update, equipment_update, etc.)
• Broadcast automatique à tous les clients
• Simulation mouvement avec positions utilisateurs

Services Externes:
• API Open-Meteo (météo temps réel)
• Géocodage automatique d'adresses
• Estimation luminosité selon météo"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(13)
        paragraph.space_after = Pt(6)
        if paragraph.text.endswith(':'):
            paragraph.font.bold = True
            paragraph.font.size = Pt(15)
            paragraph.font.color.rgb = RGBColor(37, 99, 235)
        elif paragraph.text.startswith('•'):
            paragraph.level = 1
    
    # ==================== SLIDE 10: POINTS TECHNIQUES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Points Techniques Avancés"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """Architecture Asynchrone:
• Tornado async/await natif (pas de callback hell)
• SQLAlchemy 2.0 avec asyncpg (3x plus rapide que psycopg2)
• Toutes les opérations DB non-bloquantes
• Gestion concurrence avec contextes async

Optimisations Performance:
• Tri stable des listes (évite réorganisation UI)
• Indexes sur colonnes fréquemment utilisées
• Cleanup automatique historique (limite 1000 événements)
• Pagination sur tous les endpoints de liste
• JSONB pour données flexibles (grilles maison)

Gestion Erreurs:
• Try/catch complets sur toutes les opérations async
• Validation côté serveur (email, téléphone, fichiers)
• Messages d'erreur explicites JSON
• Logging console pour debugging

Code Quality:
• Black formatter (120 caractères max)
• Flake8 linter configuré
• Type hints Python partout
• Docstrings complètes en anglais
• Architecture MVC claire (Models, Handlers, Services)

Tests & Documentation:
• 5 fichiers markdown (5000+ lignes)
• API_DOCUMENTATION.md (50+ endpoints)
• DEMONSTRATION_GUIDE.md (plan démo détaillé)
• PROJECT_REQUIREMENTS.md (conformité 35/35 points)"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.space_after = Pt(5)
        if paragraph.text.endswith(':'):
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(37, 99, 235)
        elif paragraph.text.startswith('•'):
            paragraph.level = 1
    
    # ==================== SLIDE 11: STATISTIQUES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Statistiques du Projet"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """Code Backend (Python):
• 29 fichiers Python
• 5,772 lignes de code
• 15 handlers REST API
• 9 modèles SQLAlchemy
• 3 services (météo, grille, permissions)

Code Frontend (HTML/CSS/JS):
• 21 fichiers frontend
• house.js: 2,239 lignes (logique principale)
• 10 pages HTML (SPA)
• CSS modulaire par composant

Documentation:
• 5 fichiers markdown
• 5,057 lignes de documentation
• API complète documentée
• Guide de démonstration détaillé

Base de Données:
• 9 tables normalisées
• 50+ endpoints REST API
• Relations complexes (1:N, N:M)
• Support JSONB pour flexibilité"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.space_after = Pt(10)
        if paragraph.text.endswith(':'):
            paragraph.font.bold = True
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(37, 99, 235)
        elif paragraph.text.startswith('•'):
            paragraph.level = 1
    
    # ==================== SLIDE 12: CONFORMITÉ ====================
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Conformité Cahier des Charges"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(37, 99, 235)
    
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    content = """Fonctionnalités Internes (25/25 points):
• Gestion utilisateurs: 6/6 (auth, profil, photo)
• Gestion maison: 8/8 (CRUD, pièces, capteurs, équipements, membres)
• Capteurs & équipements: 6/6 (4 types chaque, APIs dédiées)
• Interface client: 3/3 (dashboard, contrôles, météo, historique)
• Service live: 2/2 (WebSocket temps réel)

Fonctionnalités Externes (5/5 points):
• API météo Open-Meteo intégrée
• Géocodage automatique d'adresses
• Gestion erreurs complète

Système d'Information (5/5 points):
• 9 tables PostgreSQL normalisées
• Foreign keys avec contraintes
• Requêtes complexes avec jointures
• Indexes pour performance

TOTAL: 35/35 POINTS ✓

Bonus Implémentés:
• JWT pour API REST
• Middleware authentification centralisé
• Interface responsive mobile-friendly
• Éditeur graphique grille maison
• Système invitations avec notifications
• Code quality (Black, Flake8, type hints)"""
    
    tf.text = content
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.space_after = Pt(8)
        if paragraph.text.endswith(':') or 'TOTAL' in paragraph.text:
            paragraph.font.bold = True
            paragraph.font.size = Pt(16)
            if 'TOTAL' in paragraph.text:
                paragraph.font.color.rgb = RGBColor(34, 197, 94)
            else:
                paragraph.font.color.rgb = RGBColor(37, 99, 235)
        elif paragraph.text.startswith('•'):
            paragraph.level = 1
    
    # ==================== SLIDE 13: CONCLUSION ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusion"
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(37, 99, 235)
    p.alignment = PP_ALIGN.CENTER
    
    # Points clés
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(7), Inches(3))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content = """Projet complet et fonctionnel:
✓ Architecture moderne et scalable
✓ Code quality élevée (formaté, documenté, typé)
✓ Performance optimisée (async, indexes, pagination)
✓ Sécurité robuste (bcrypt, validation, permissions)
✓ Interface utilisateur intuitive et responsive
✓ Documentation exhaustive (5000+ lignes)

Technologies maîtrisées:
Python • Tornado • PostgreSQL • SQLAlchemy
JavaScript • WebSocket • REST API • Async/Await

GitHub: github.com/Devid-yl/smarthome"""
    
    content_frame.text = content
    
    for i, paragraph in enumerate(content_frame.paragraphs):
        paragraph.font.size = Pt(16)
        paragraph.space_after = Pt(8)
        paragraph.alignment = PP_ALIGN.LEFT
        
        if i == 0 or i == 7:  # Titres
            paragraph.font.bold = True
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(37, 99, 235)
        elif paragraph.text.startswith('✓'):
            paragraph.font.color.rgb = RGBColor(34, 197, 94)
        elif i == 10:  # GitHub
            paragraph.font.color.rgb = RGBColor(100, 116, 139)
            paragraph.font.italic = True
    
    # Merci
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Merci pour votre attention !"
    
    p = thanks_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(100, 116, 139)
    p.alignment = PP_ALIGN.CENTER
    
    # Sauvegarder
    filename = 'SmartHome_Presentation.pptx'
    prs.save(filename)
    print(f"✅ Présentation créée: {filename}")
    print(f"📊 {len(prs.slides)} slides générées")
    print(f"⏱️  Durée estimée: 7-10 minutes")

if __name__ == '__main__':
    create_presentation()
